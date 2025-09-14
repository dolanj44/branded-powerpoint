from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_ai_workflow_presentation():
    # Create presentation
    prs = Presentation()
    
    # Define colors for consistent branding (professional blue theme)
    primary_blue = RGBColor(31, 73, 125)
    accent_blue = RGBColor(79, 129, 189)
    light_gray = RGBColor(242, 242, 242)
    dark_gray = RGBColor(64, 64, 64)
    
    # Slide 1: Title Slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    title1 = slide1.shapes.title
    subtitle1 = slide1.placeholders[1]
    
    title1.text = "AI-Enhanced Workflow Proposal: Weekly Status Reporting"
    title1.text_frame.paragraphs[0].font.color.rgb = primary_blue
    title1.text_frame.paragraphs[0].font.size = Pt(44)
    
    subtitle1.text = "Driving efficiency, consistency, and insight across 11 project teams"
    subtitle1.text_frame.paragraphs[0].font.color.rgb = dark_gray
    subtitle1.text_frame.paragraphs[0].font.size = Pt(20)
    
    # Slide 2: Business Need
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    title2 = slide2.shapes.title
    content2 = slide2.placeholders[1]
    
    title2.text = "Business Need"
    title2.text_frame.paragraphs[0].font.color.rgb = primary_blue
    
    content2.text = """• Weekly status reporting across 11 teams is highly manual and inconsistent
• Current process consumes ~20 hours weekly across the program
• Leadership requires timely, consistent, accurate insights to make informed decisions
• Growing volume of projects demands better efficiency"""
    
    for paragraph in content2.text_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = dark_gray
    
    # Slide 3: Current Process (Pain Points)
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    title3 = slide3.shapes.title
    content3 = slide3.placeholders[1]
    
    title3.text = "Current Process Pain Points"
    title3.text_frame.paragraphs[0].font.color.rgb = primary_blue
    
    content3.text = """• Manual data collection from Jira, Slack, Confluence, SharePoint
• Drafting & polishing reports takes ~2 hours per team
• Inconsistent outputs, variation across teams
• Delays in surfacing risks and dependencies
• Repetitive effort leads to missed themes across projects
• Valuable team time lost to formatting instead of strategic work"""
    
    for paragraph in content3.text_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = dark_gray
    
    # Slide 4: Proposed Solution
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    title4 = slide4.shapes.title
    content4 = slide4.placeholders[1]
    
    title4.text = "Proposed AI Solution"
    title4.text_frame.paragraphs[0].font.color.rgb = primary_blue
    
    content4.text = """• RAG models grounded in Jira, Confluence, and SharePoint
• AI agents scan Slack channels for highlights and auto-generate charts
• Intelligent prompting for drafting, structuring, and polishing reports
• Consistent executive-ready format across all teams
• AI provides first draft; team leads validate and finalize
• Real-time project updates and documentation integration"""
    
    for paragraph in content4.text_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = dark_gray
    
    # Slide 5: Business Impact - Time Savings
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    title5 = slide5.shapes.title
    content5 = slide5.placeholders[1]
    
    title5.text = "Business Impact: Time Savings"
    title5.text_frame.paragraphs[0].font.color.rgb = primary_blue
    
    content5.text = """Time Savings per Team:
• Current: 110 minutes per week
• With AI: 37 minutes per week
• Net savings: 73 minutes per team

Program-wide Impact:
• 13.5 hours saved weekly across 11 teams
• ~160 hours saved per PI cycle
• 66% reduction in report preparation time"""
    
    for paragraph in content5.text_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = dark_gray
    
    # Slide 6: Quality & Strategic Gains
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    title6 = slide6.shapes.title
    content6 = slide6.placeholders[1]
    
    title6.text = "Quality & Strategic Value"
    title6.text_frame.paragraphs[0].font.color.rgb = primary_blue
    
    content6.text = """Quality Gains:
• Consistent structure & tone across all reports
• Better visibility into cross-team risks/dependencies
• Automated summarization reduces human error
• More accurate, reliable charts and metrics

Strategic Value:
• Program-wide reporting alignment
• Faster, more informed leadership decisions
• Digital transformation alignment
• Capacity shift to higher-value strategic work"""
    
    for paragraph in content6.text_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = dark_gray
    
    # Slide 7: Security Considerations
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    title7 = slide7.shapes.title
    content7 = slide7.placeholders[1]
    
    title7.text = "Security Considerations"
    title7.text_frame.paragraphs[0].font.color.rgb = primary_blue
    
    content7.text = """• Data sources: Jira, Confluence, SharePoint, Slack
• Sensitive data remains within secure, access-controlled environments
• AI tools must comply with organizational security policies
• Full compliance with privacy and audit requirements
• Secure authentication and data handling protocols
• Regular security reviews and monitoring"""
    
    for paragraph in content7.text_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = dark_gray
    
    # Slide 8: Implementation Support
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    title8 = slide8.shapes.title
    content8 = slide8.placeholders[1]
    
    title8.text = "Implementation Support Available"
    title8.text_frame.paragraphs[0].font.color.rgb = primary_blue
    
    content8.text = """• Subject matter expertise on current workflows
• Partnership on pilot testing and validation
• Development of documentation and user guides
• Collection and relay of user feedback for improvements
• Training and onboarding support for team leads
• Ongoing optimization and enhancement"""
    
    for paragraph in content8.text_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = dark_gray
    
    # Slide 9: Timeline & Next Steps
    slide9 = prs.slides.add_slide(prs.slide_layouts[1])
    title9 = slide9.shapes.title
    content9 = slide9.placeholders[1]
    
    title9.text = "Timeline & Next Steps"
    title9.text_frame.paragraphs[0].font.color.rgb = primary_blue
    
    content9.text = """Target Timeline:
• Pilot implementation within next PI (3 months)
• Aligns with leadership request for timely reporting
• Supports ongoing digital transformation goals

Next Steps:
• Secure resource approval for AI engineering expertise
• Design pilot with 1-2 teams for validation
• Measure impact: time savings, consistency, adoption
• Scale to all 11 teams based on pilot results"""
    
    for paragraph in content9.text_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = dark_gray
    
    # Slide 10: Executive Summary
    slide10 = prs.slides.add_slide(prs.slide_layouts[1])
    title10 = slide10.shapes.title
    content10 = slide10.placeholders[1]
    
    title10.text = "Executive Summary"
    title10.text_frame.paragraphs[0].font.color.rgb = primary_blue
    
    content10.text = """Key Benefits:
• ~160 hours saved per PI cycle
• Consistent reporting across 11 teams
• Faster, more informed leadership decisions
• Enhanced visibility into cross-team dependencies

Requirements:
• AI engineering expertise for secure implementation
• Integration with existing enterprise tools
• Compliance with security and privacy policies

ROI: Significant time savings + improved decision-making capability"""
    
    for paragraph in content10.text_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = dark_gray
    
    # Save presentation
    prs.save('/Users/josephdolan/projects/workflow_claude code/AI_Workflow_Proposal.pptx')
    print("PowerPoint presentation created successfully: AI_Workflow_Proposal.pptx")

if __name__ == "__main__":
    create_ai_workflow_presentation()