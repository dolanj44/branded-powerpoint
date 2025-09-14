from pptx import Presentation

def add_text_to_slide(slide, text, exclude_title=True):
    """Add text to slide by finding existing text frames or placeholders"""
    
    # First try to find text frames in shapes (excluding title if requested)
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame') and shape.text_frame:
            # Skip title shape if exclude_title is True
            if exclude_title and shape == slide.shapes.title:
                continue
            
            # Check if this text frame is empty or has placeholder text
            if not shape.text_frame.text or shape.text_frame.text.strip() in ['', 'Click to add text']:
                shape.text_frame.text = text
                return True
    
    # If no existing text frames found, try placeholders
    for placeholder in slide.placeholders:
        if hasattr(placeholder, 'text_frame') and placeholder.text_frame:
            # Skip title placeholder if exclude_title is True
            if exclude_title and placeholder == slide.shapes.title:
                continue
                
            if not placeholder.text or placeholder.text.strip() in ['', 'Click to add text']:
                placeholder.text = text
                return True
    
    return False

def add_text_to_multiple_areas(slide, texts, exclude_title=True):
    """Add multiple text blocks to different areas of a slide"""
    
    # Collect all available text containers (excluding title if requested)
    text_containers = []
    
    # First collect text frames from shapes
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame') and shape.text_frame:
            if exclude_title and shape == slide.shapes.title:
                continue
            text_containers.append(shape)
    
    # Then collect placeholders
    for placeholder in slide.placeholders:
        if hasattr(placeholder, 'text_frame') and placeholder.text_frame:
            if exclude_title and placeholder == slide.shapes.title:
                continue
            # Only add if not already in shapes list
            if placeholder not in [shape for shape in slide.shapes]:
                text_containers.append(placeholder)
    
    # Add texts to containers
    for i, text in enumerate(texts):
        if i < len(text_containers):
            if hasattr(text_containers[i], 'text'):
                text_containers[i].text = text
            elif hasattr(text_containers[i], 'text_frame'):
                text_containers[i].text_frame.text = text

def create_presentation_with_proper_text_handling():
    template_path = 'alot going on.pptx'
    
    try:
        print("Loading template and analyzing existing text frames...")
        prs = Presentation(template_path)
        
        # Clear existing slides but keep all slide masters
        slide_count = len(prs.slides)
        for i in range(slide_count - 1, -1, -1):
            r_id = prs.slides._sldIdLst[i].rId
            prs.part.drop_rel(r_id)
            del prs.slides._sldIdLst[i]
        
        print(f"Cleared {slide_count} existing slides")
        
        # Access layouts from specific slide masters
        title_master = prs.slide_masters[0]  # Title slides
        content_master = prs.slide_masters[1]  # Content slides  
        summary_master = prs.slide_masters[2]  # Summary slides
        divider_master = prs.slide_masters[3]  # Divider slides
        
        # Get specific layouts
        title_color_layout = title_master.slide_layouts[0]  # Title Slide - Color
        
        content_one_col = content_master.slide_layouts[0]  # Content - One Column
        content_two_col = content_master.slide_layouts[1]  # Content - Two Column
        content_three_col = content_master.slide_layouts[2]  # Content - Three Column
        content_half_image = content_master.slide_layouts[4]  # Content - Half Image
        
        summary_exec = summary_master.slide_layouts[1]  # Summary - Exec
        
        divider_texture1 = divider_master.slide_layouts[0]  # Divider - Texture 1
        
        print("Creating slides with proper text handling...")
        
        # Slide 1: Title Slide
        slide = prs.slides.add_slide(title_color_layout)
        slide.shapes.title.text = "AI-Enhanced Workflow Proposal"
        
        # Add subtitle to existing text frame
        subtitle_text = "Weekly Status Reporting: Driving efficiency across 11 project teams"
        add_text_to_slide(slide, subtitle_text, exclude_title=True)
        
        # Slide 2: Section Divider
        slide = prs.slides.add_slide(divider_texture1)
        slide.shapes.title.text = "Current State Analysis"
        
        # Slide 3: Business Need
        slide = prs.slides.add_slide(content_one_col)
        slide.shapes.title.text = "Business Need"
        
        business_need_text = """• Weekly status reporting across 11 teams is highly manual and inconsistent
• Current process consumes ~20 hours weekly across the program  
• Leadership requires timely, consistent, accurate insights to make informed decisions
• Growing volume of projects demands better efficiency and standardization"""
        
        add_text_to_slide(slide, business_need_text, exclude_title=True)
        
        # Slide 4: Current Pain Points (Two Column)
        slide = prs.slides.add_slide(content_two_col)
        slide.shapes.title.text = "Current Process Pain Points"
        
        pain_points_texts = [
            """Manual Processes:
• Data collection from Jira, Slack, Confluence, SharePoint
• Drafting & polishing takes ~2 hours per team
• Repetitive formatting and structuring work""",
            
            """Quality Issues:
• Inconsistent outputs across teams
• Delays in surfacing risks and dependencies  
• Missed themes and strategic insights"""
        ]
        
        add_text_to_multiple_areas(slide, pain_points_texts, exclude_title=True)
        
        # Slide 5: Section Divider
        slide = prs.slides.add_slide(divider_texture1)
        slide.shapes.title.text = "Proposed AI Solution"
        
        # Slide 6: AI Solution Overview
        slide = prs.slides.add_slide(content_half_image)
        slide.shapes.title.text = "AI-Enhanced Workflow Components"
        
        ai_solution_text = """• RAG models grounded in Jira, Confluence, and SharePoint
• AI agents scan Slack channels for highlights  
• Auto-generate charts and visual insights
• Intelligent prompting for consistent executive formatting
• Real-time project updates and documentation integration
• Human validation of AI-generated first drafts"""
        
        add_text_to_slide(slide, ai_solution_text, exclude_title=True)
        
        # Slide 7: Time Savings (Three Column)
        slide = prs.slides.add_slide(content_three_col)
        slide.shapes.title.text = "Time Savings Impact"
        
        time_savings_texts = [
            """Current State
110 minutes per team
per week""",
            
            """With AI
37 minutes per team
per week""",
            
            """Net Savings
73 minutes per team
66% reduction"""
        ]
        
        add_text_to_multiple_areas(slide, time_savings_texts, exclude_title=True)
        
        # Slide 8: Program Impact
        slide = prs.slides.add_slide(content_one_col)
        slide.shapes.title.text = "Program-Wide Impact"
        
        program_impact_text = """Time Savings:
• 13.5 hours saved weekly across 11 teams
• ~160 hours saved per PI cycle
• Equivalent to 4 full work weeks of capacity returned

Quality Improvements:
• Consistent structure and tone across all reports
• Enhanced visibility into cross-team themes and dependencies  
• Automated summarization reduces human error
• More accurate, reliable charts and metrics"""
        
        add_text_to_slide(slide, program_impact_text, exclude_title=True)
        
        # Slide 9: Strategic Value (Two Column)
        slide = prs.slides.add_slide(content_two_col)
        slide.shapes.title.text = "Strategic & Quality Benefits"
        
        strategic_texts = [
            """Strategic Value:
• Program-wide reporting alignment
• Faster, more informed leadership decisions
• Digital transformation alignment with measurable impact
• Capacity shift: frees team leads for strategic work""",
            
            """Security & Compliance:
• Data remains in secure, access-controlled environments
• Full compliance with organizational policies
• Regular security reviews and audit trails
• Enterprise data governance alignment"""
        ]
        
        add_text_to_multiple_areas(slide, strategic_texts, exclude_title=True)
        
        # Slide 10: Implementation
        slide = prs.slides.add_slide(content_one_col)
        slide.shapes.title.text = "Implementation Support & Timeline"
        
        implementation_text = """Available Support:
• Subject matter expertise on current workflows and pain points
• Partnership on pilot testing and validation of AI-generated reports  
• Development of comprehensive documentation and user guides
• Training and onboarding support for team leads

Timeline: Pilot within next PI (3 months)
• Aligns with leadership's request for timely reporting
• Design pilot with 1-2 teams → validate → scale to all 11 teams"""
        
        add_text_to_slide(slide, implementation_text, exclude_title=True)
        
        # Slide 11: Executive Summary
        slide = prs.slides.add_slide(summary_exec)
        slide.shapes.title.text = "Executive Summary"
        
        exec_summary_texts = [
            "~160 hours saved per PI cycle",
            "Consistent reporting across 11 teams",  
            "Enhanced leadership decision-making",
            "Measurable AI transformation progress",
            "Requires AI engineering expertise + secure implementation"
        ]
        
        add_text_to_multiple_areas(slide, exec_summary_texts, exclude_title=True)
        
        # Save the presentation
        output_path = 'AI_Workflow_Proposal_Final.pptx'
        prs.save(output_path)
        print(f"\nFinal presentation created: {output_path}")
        print(f"Total slides: {len(prs.slides)}")
        
        # Print summary of layouts used
        print("\nLayouts used:")
        for i, slide in enumerate(prs.slides):
            print(f"  Slide {i+1}: {slide.slide_layout.name}")
        
        return output_path
        
    except Exception as e:
        print(f"Error creating presentation: {e}")
        import traceback
        traceback.print_exc()
        return None

if __name__ == "__main__":
    create_presentation_with_proper_text_handling()